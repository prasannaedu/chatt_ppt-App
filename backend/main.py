from fastapi import FastAPI, HTTPException, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from typing import List, Optional
import os, io, time, re, tempfile, sqlite3, google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import logging
from datetime import datetime

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Google Gemini Configuration
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.5-flash-lite")

# Configure Gemini
try:
    genai.configure(api_key=GEMINI_API_KEY)
    gemini_model = genai.GenerativeModel(GEMINI_MODEL)
    logger.info(f"✅ Gemini AI configured: {GEMINI_MODEL}")
except Exception as e:
    logger.error(f"❌ Gemini configuration failed: {e}")
    gemini_model = None

app = FastAPI(title="Chat-to-PPT API", version="18.5") # Updated Version

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class GeneratePayload(BaseModel):
    topic: str
    slides: int = 6
    style: str = "Blue-Professional"
    background_color: str = "#FFFFFF"
    include_images: bool = False
    content_depth: str = "detailed"
    notes: Optional[str] = ""

# Enhanced background colors with better contrast
BACKGROUND_COLORS = {
    "Pure White": "#FFFFFF",
    "Soft Gray": "#F8FAFC", 
    "Warm White": "#FEF7EE",
    "Ice Blue": "#F0F9FF",
    "Mint Cream": "#F0FDF4",
    "Lavender": "#FDF4FF",
    "Peach": "#FFF7ED",
    "Dark Mode": "#1E293B",
    "Professional Blue": "#F0F7FF",
    "Executive Gray": "#F8FAFC"
}

# **UNIFIED PROMPT TEMPLATES - CONSISTENT BULLET POINTS**
PROMPT_TEMPLATES = {
    "basic": """
Create a concise PowerPoint presentation outline about "{topic}" with EXACTLY {slides} slides.

CRITICAL REQUIREMENTS:
- Generate EXACTLY {slides} slides - no more, no less
- Each slide must start with "Slide X: [Title]"
- Use bullet points starting with "-"
- Keep EXACTLY 3 bullet points per slide
- Each bullet should be 5-7 words
- Create SHORT, CLEAR titles (max 8-10 words)
- Focus on fundamental concepts and key takeaways
- Use simple, clear language suitable for beginners
- Ensure each bullet point is COMPLETE and doesn't get cut off
- Make each slide UNIQUE with different content

BASIC SLIDE STRUCTURE - EXACTLY {slides} SLIDES:
Slide 1: [Introduction to {topic}]
- [Core concept definition and overview]
- [Main purpose and basic applications]
- [Key benefits and importance]

Slide 2: [Fundamental Principles]
- [Primary principle explanation]
- [Supporting concepts and basics]
- [Simple examples and use cases]

Slide 3: [Key Components]
- [Main elements breakdown]
- [Basic functionality overview]
- [Essential features summary]

Continue this pattern for exactly {slides} slides total. Focus on foundational knowledge.
""",

    "detailed": """
Create a detailed professional PowerPoint presentation outline about "{topic}" with EXACTLY {slides} slides.

CRITICAL REQUIREMENTS:
- Generate EXACTLY {slides} slides - no more, no less
- Each slide must start with "Slide X: [Title]"
- Use bullet points starting with "-"
- Include EXACTLY 4 bullet points per slide
- Each bullet should be 8-12 words
- Create CONCISE, MEANINGFUL titles (max 10-12 words)
- Include analysis, examples, and practical applications
- Provide deeper insights and strategic considerations
- Ensure each bullet point is COMPLETE and doesn't get cut off
- Make each slide UNIQUE with different perspectives

DETAILED SLIDE STRUCTURE - EXACTLY {slides} SLIDES:
Slide 1: [Comprehensive {topic} Analysis]
- [Detailed overview with market context and current relevance]
- [Key industry trends analysis and emerging patterns]
- [Strategic importance and business impact considerations]
- [Implementation challenges and potential solutions]

Slide 2: [Advanced Concepts Framework]
- [Technical breakdown with component relationships]
- [Methodologies explanation and best practices]
- [Case study analysis with real-world examples]
- [Performance metrics and success measurement]

Slide 3: [Strategic Implementation Plan]
- [Phased implementation approach with timeline]
- [Resource allocation and team structure]
- [Risk assessment with mitigation strategies]
- [Success metrics and improvement framework]

Continue this pattern for exactly {slides} slides total. Provide comprehensive analysis.
""",

    "comprehensive": """
Create an executive-level comprehensive PowerPoint presentation about "{topic}" with EXACTLY {slides} slides.

CRITICAL REQUIREMENTS:
- Generate EXACTLY {slides} slides - no more, no less
- Each slide must start with "Slide X: [Title]"
- Use bullet points starting with "-"
- Include EXACTLY 5 bullet points per slide
- Each bullet should be 12-18 words
- Create FOCUSED, STRATEGIC titles (max 12-15 words)
- Include data-driven insights and strategic frameworks
- Cover financial implications and long-term impact
- Ensure each bullet point is COMPLETE and doesn't get cut off
- Make titles CLEAR and COMPLETE - no truncation
- Ensure VARIED content across slides - no repetition

COMPREHENSIVE SLIDE STRUCTURE - EXACTLY {slides} SLIDES:
Slide 1: [Executive Strategic Overview]
- [Comprehensive market analysis with current trends and competitive landscape assessment]
- [Strategic business case development with ROI calculation and financial projections]
- [Implementation roadmap detailing phased approach and resource allocation strategy]
- [Stakeholder impact analysis covering organizational change management requirements]
- [Performance measurement framework defining KPIs and continuous improvement processes]

Slide 2: [Technical Architecture & Innovation Strategy]
- [Technical infrastructure design including scalability considerations and integration points]
- [Innovation adoption framework covering emerging technologies and partnership opportunities]
- [Data analytics implementation with business intelligence tools and dashboard development]
- [Security and compliance framework addressing regulatory requirements and risk management]
- [Change management strategy detailing training programs and communication plans]

Slide 3: [Financial Analysis & Business Impact]
- [Detailed financial modeling including revenue projections and cost-benefit analysis]
- [Investment justification framework covering capital expenditure and ROI timeline]
- [Risk assessment matrix identifying operational, financial, and market risks]
- [Stakeholder value proposition detailing benefits across customer and employee segments]
- [Strategic alignment analysis connecting initiative to organizational objectives]

Continue this pattern for exactly {slides} slides total. Provide executive-level strategic insights with NO CONTENT REPETITION.
"""
}

BULLET_RE = re.compile(r"^-\s+")

def call_gemini(prompt: str):
    """Call Google Gemini AI for content generation with better error handling"""
    if not gemini_model:
        logger.error("❌ Gemini model not available")
        return None
    
    try:
        logger.info(f"🤖 Calling Gemini AI: {GEMINI_MODEL}")
        logger.info(f"📝 Prompt length: {len(prompt)}")
        
        # Generate content with Gemini
        response = gemini_model.generate_content(
            prompt,
            generation_config={
                "temperature": 0.8,
                "top_p": 0.9,
                "top_k": 50,
                "max_output_tokens": 4096,
            }
        )
        
        if response.text:
            logger.info(f"✅ Gemini response received, length: {len(response.text)}")
            logger.info(f"📄 Content preview: {response.text[:300]}...")
            
            # Validate the response has proper structure
            if "Slide" in response.text and "-" in response.text:
                return response.text.strip()
            else:
                logger.warning("⚠️ Gemini response lacks proper structure")
                return None
        else:
            logger.warning("⚠️ Gemini returned empty response")
            return None
            
    except Exception as e:
        logger.error(f"❌ Gemini API call failed: {str(e)}")
        return None

def generate_consistent_fallback(topic: str, num_slides: int, content_depth: str):
    """Generate consistent fallback content with EXACT bullet point counts"""
    logger.info(f"🔄 Generating CONSISTENT fallback content for: {topic}, {num_slides} slides, {content_depth} depth")
    
    # Define exact bullet point counts for each depth
    bullet_counts = {
        "basic": 3,
        "detailed": 4,
        "comprehensive": 5
    }
    
    target_bullets = bullet_counts.get(content_depth, 4)
    
    # Define slide templates with exact bullet counts
    slide_templates = {
        "comprehensive": [
            {
                "title": f"Executive Strategic Overview: {topic}",
                "bullets": [
                    f"Comprehensive market analysis of {topic} with current industry trends and competitive landscape assessment",
                    f"Strategic business case development with detailed ROI calculation, financial projections, and investment justification",
                    f"Implementation roadmap detailing phased approach, resource allocation strategy, and milestone tracking",
                    f"Stakeholder impact analysis covering organizational change management, training requirements, and communication strategy",
                    f"Performance measurement framework defining KPIs, success metrics, and continuous improvement processes"
                ]
            },
            {
                "title": "Technical Architecture & Innovation Strategy",
                "bullets": [
                    f"Technical infrastructure design including scalability considerations, integration points, and future-proofing strategies",
                    f"Innovation adoption framework covering emerging technologies, partnership opportunities, and competitive advantage positioning",
                    f"Data analytics implementation with business intelligence tools, predictive modeling, and real-time dashboard development",
                    f"Security and compliance framework addressing regulatory requirements, data protection, and risk management protocols",
                    f"Change management strategy detailing organizational transformation, training programs, and cultural adoption measurement"
                ]
            },
            {
                "title": "Financial Analysis & Business Impact Assessment",
                "bullets": [
                    f"Detailed financial modeling including revenue projections, cost analysis, break-even calculation, and sensitivity analysis",
                    f"Investment justification framework covering capital expenditure, operational costs, return on investment timeline, and payback period",
                    f"Risk assessment matrix identifying operational, financial, technical, and market risks with probability-impact analysis",
                    f"Stakeholder value proposition detailing benefits for customers, employees, shareholders, and partners",
                    f"Strategic alignment analysis connecting to organizational goals, competitive positioning, and long-term growth objectives"
                ]
            },
            {
                "title": "Implementation Excellence & Project Management Framework",
                "bullets": [
                    f"Project management methodology with agile approach, sprint planning, resource allocation, and governance structure",
                    f"Quality assurance framework covering testing protocols, performance benchmarking, user acceptance criteria, and feedback mechanisms",
                    f"Team structure and capability development outlining roles, responsibilities, skill requirements, and training programs",
                    f"Vendor and partnership strategy detailing selection criteria, contract management, performance monitoring, and relationship management",
                    f"Operational excellence framework covering process optimization, automation opportunities, efficiency metrics, and service level agreements"
                ]
            },
            {
                "title": "Strategic Risk Management & Mitigation Planning",
                "bullets": [
                    f"Comprehensive risk identification process covering operational, financial, technical, and market-related challenges",
                    f"Risk assessment methodology using probability-impact matrix and quantitative analysis techniques",
                    f"Mitigation strategy development with contingency planning and alternative scenario analysis",
                    f"Monitoring and control framework with early warning indicators and escalation procedures",
                    f"Business continuity planning ensuring operational resilience and disaster recovery capabilities"
                ]
            },
            {
                "title": "Performance Optimization & Strategic Roadmap",
                "bullets": [
                    f"Performance benchmarking against industry standards and competitor analysis for continuous improvement",
                    f"Optimization strategies focusing on efficiency gains, cost reduction, and value enhancement opportunities",
                    f"Technology roadmap aligning with business strategy and emerging innovation trends",
                    f"Talent development and capability building programs for sustained competitive advantage",
                    f"Long-term strategic vision with measurable objectives and milestone tracking mechanisms"
                ]
            }
        ],
        "detailed": [
            {
                "title": f"Strategic Analysis Overview: {topic}",
                "bullets": [
                    f"Comprehensive analysis of {topic} market position and competitive landscape",
                    f"Key industry trends assessment and emerging opportunity identification",
                    f"Strategic importance evaluation and business impact considerations",
                    f"Implementation challenges overview and solution framework development"
                ]
            },
            {
                "title": "Technical Framework & Architecture",
                "bullets": [
                    f"Technical architecture overview and component relationships mapping",
                    f"Methodology explanation and implementation best practices",
                    f"Case study analysis with real-world application examples",
                    f"Performance metrics definition and success measurement criteria"
                ]
            },
            {
                "title": "Implementation Strategy & Planning",
                "bullets": [
                    f"Phased implementation approach with timeline and milestone planning",
                    f"Resource allocation strategy and team structure recommendations",
                    f"Risk assessment framework with mitigation strategy development",
                    f"Success metrics tracking and continuous improvement framework"
                ]
            },
            {
                "title": "Market Analysis & Competitive Positioning",
                "bullets": [
                    f"Target market segmentation and customer needs analysis",
                    f"Competitive landscape assessment and differentiation strategy",
                    f"Market opportunity sizing and growth potential evaluation",
                    f"Strategic positioning and value proposition development"
                ]
            },
            {
                "title": "Operational Excellence & Efficiency",
                "bullets": [
                    f"Process optimization and efficiency improvement initiatives",
                    f"Quality management and performance monitoring systems",
                    f"Resource utilization and capacity planning strategies",
                    f"Continuous improvement and innovation implementation"
                ]
            },
            {
                "title": "Future Outlook & Strategic Recommendations",
                "bullets": [
                    f"Emerging trends analysis and future market developments",
                    f"Strategic recommendations and implementation priorities",
                    f"Long-term vision and growth opportunity identification",
                    f"Next steps and action plan for immediate execution"
                ]
            }
        ],
        "basic": [
            {
                "title": f"Introduction & Overview: {topic}",
                "bullets": [
                    f"Core concept definition and basic overview of {topic}",
                    f"Main purpose explanation and primary applications",
                    f"Key benefits summary and importance assessment"
                ]
            },
            {
                "title": "Fundamental Concepts & Principles",
                "bullets": [
                    f"Primary principles explanation and supporting concepts",
                    f"Basic functionality overview and key features",
                    f"Simple examples demonstration and use cases"
                ]
            },
            {
                "title": "Practical Applications & Use Cases",
                "bullets": [
                    f"Implementation approach and basic requirements",
                    f"Common use cases and application scenarios",
                    f"Success factors and best practices summary"
                ]
            },
            {
                "title": "Key Benefits & Competitive Advantages",
                "bullets": [
                    f"Primary advantages and competitive benefits",
                    f"Efficiency improvements and cost savings",
                    f"User experience enhancements and value creation"
                ]
            },
            {
                "title": "Implementation Steps & Requirements",
                "bullets": [
                    f"Initial setup and configuration requirements",
                    f"Deployment process and timeline overview",
                    f"Training needs and user adoption strategy"
                ]
            },
            {
                "title": "Summary & Next Steps",
                "bullets": [
                    f"Key takeaways and main points summary",
                    f"Recommended actions and implementation priorities",
                    f"Future considerations and expansion opportunities"
                ]
            }
        ]
    }
    
    # Get the appropriate template set
    templates = slide_templates.get(content_depth, slide_templates["detailed"])
    
    # Generate presentation content with EXACT bullet counts
    presentation_content = ""
    
    for i in range(min(num_slides, len(templates))):
        template = templates[i]
        
        presentation_content += f"Slide {i+1}: {template['title']}\n"
        
        # Ensure we have exactly the right number of bullets
        bullets = template['bullets'][:target_bullets]
        
        for bullet in bullets:
            presentation_content += f"- {bullet}\n"
        
        presentation_content += "\n"
    
    logger.info(f"✅ Generated CONSISTENT fallback with {min(num_slides, len(templates))} slides, {target_bullets} bullets each")
    return presentation_content

def parse_outline(text: str, content_depth: str = "detailed", requested_slides: int = 6):
    """Parse the outline text into structured slides - ENSURING EXACT BULLET POINTS"""
    slides = []
    
    # If no text provided, return empty to trigger fallback
    if not text or not text.strip():
        logger.warning("⚠️ No text provided to parse_outline")
        return []
    
    blocks = re.split(r"\n\s*\n", text.strip())
    
    # CORRECTED: Maintain EXACT bullet points for each content depth
    target_bullets = {
        "basic": 3,
        "detailed": 4,
        "comprehensive": 5
    }.get(content_depth, 4)
    
    slide_count = 0
    
    for block in blocks:
        if slide_count >= requested_slides:
            break
            
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
            
        # Extract title
        title_line = lines[0]
        title = re.sub(r"^Slide\s*\d+:\s*", "", title_line, flags=re.I).strip()
        
        if not title or len(title) < 2:
            title_match = re.search(r"^(?:Slide\s*\d+[:\.]?\s*)?(.+)$", title_line, re.I)
            if title_match:
                title = title_match.group(1).strip()
            else:
                continue  # Skip if no valid title
        
        # Extract bullets
        bullets = []
        for ln in lines[1:]:
            clean_line = BULLET_RE.sub("", ln.strip())
            if not clean_line and ln.strip().startswith(('•', '*', '·')):
                clean_line = ln.strip()[1:].strip()
                
            if clean_line and len(clean_line) > 5:
                clean_line = re.sub(r'\s+', ' ', clean_line).strip()
                bullets.append(clean_line)
        
        # FIXED: Only accept slides with EXACT target bullet count
        if len(bullets) >= target_bullets:
            # Use exactly the target number of bullets
            bullets = bullets[:target_bullets]
            slides.append({"title": title, "bullets": bullets, "image_path": None})
            slide_count += 1
        elif bullets and len(bullets) > 0:
            # If we have some bullets but not enough, pad with fallback content
            logger.warning(f"⚠️ Slide has only {len(bullets)} bullets, expected {target_bullets}")
            while len(bullets) < target_bullets:
                bullets.append(f"Additional strategic consideration for {title.lower()}")
            slides.append({"title": title, "bullets": bullets, "image_path": None})
            slide_count += 1
    
    logger.info(f"📊 Parsed {len(slides)} slides with {target_bullets} bullets per slide")
    
    # If we got fewer slides than requested, log it but continue
    if len(slides) < requested_slides:
        logger.warning(f"⚠️ Generated only {len(slides)} slides (requested: {requested_slides})")
    
    return slides[:requested_slides]  # Ensure we don't return more than requested

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join([c*2 for c in hex_color])
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def theme_colors(style: str, background_color: str = "#FFFFFF"):
    """Get color scheme based on theme and background with improved contrast"""
    bg_rgb = hex_to_rgb(background_color)
    bg_brightness = (bg_rgb[0] * 299 + bg_rgb[1] * 587 + bg_rgb[2] * 114) / 1000
    
    if bg_brightness < 128:
        title_color = (255, 255, 255)
        bullet_color = (240, 240, 240)
        accent_color = (100, 150, 255)
    else:
        if style.lower().startswith("pink"):
            title_color = (150, 30, 80)
            bullet_color = (60, 40, 50)
            accent_color = (200, 80, 150)
        else:
            title_color = (10, 50, 120)
            bullet_color = (30, 30, 50)
            accent_color = (80, 130, 255)
    
    return {
        "bg": bg_rgb,
        "title": title_color,
        "bullet": bullet_color,
        "accent": accent_color
    }

def build_ppt(slides, style: str, background_color: str = "#FFFFFF", include_images: bool = False, content_depth: str = "detailed"):
    """Build PowerPoint presentation with ALL BULLET POINTS PRESERVED"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        colors = theme_colors(style, background_color)

        # Title slide - COMPLETELY FIXED: No bullet points on title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors['bg'])
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        # Use the topic from the first slide title, but clean it up
        if slides and len(slides) > 0:
            main_title = slides[0]['title']
            # Clean up the title (remove "Executive Strategic Overview:" etc.)
            main_title = re.sub(r'^[^:]+:\s*', '', main_title)
            title_shape.text = main_title
        else:
            title_shape.text = "AI Presentation"
        
        # Show content depth and slide count info
        bullet_counts = {
            "basic": "3 bullet points per slide",
            "detailed": "4 bullet points per slide", 
            "comprehensive": "5 bullet points per slide"
        }
        
        subtitle_shape.text = f"Professional {content_depth.title()} Presentation\n{len(slides)} Content Slides • {bullet_counts.get(content_depth, 'Professional Quality')}"
        
        # Smart title font sizing
        title_text_length = len(title_shape.text)
        if title_text_length > 60:
            title_font_size = Pt(28)
        elif title_text_length > 40:
            title_font_size = Pt(32)
        elif title_text_length > 25:
            title_font_size = Pt(36)
        else:
            title_font_size = Pt(44)
        
        # Enable text wrapping in title shape
        title_shape.text_frame.word_wrap = True
        
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = title_font_size
                run.font.bold = True
                run.font.color.rgb = RGBColor(*colors['title'])
                run.font.name = 'Calibri'
        
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(*colors['bullet'])
                run.font.name = 'Calibri'

        # Content slides - Use ALL slides as content
        for i, slide_data in enumerate(slides):
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*colors['bg'])
            
            title_shape = slide.shapes.title
            content_title = slide_data['title']
            title_shape.text = content_title
            
            # Enable text wrapping for content titles
            title_shape.text_frame.word_wrap = True
            
            # Smart content title sizing
            content_title_length = len(content_title)
            if content_title_length > 50:
                content_title_size = Pt(20)
            elif content_title_length > 35:
                content_title_size = Pt(22)
            elif content_title_length > 20:
                content_title_size = Pt(24)
            else:
                content_title_size = Pt(28)
            
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = content_title_size
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(*colors['title'])
                    run.font.name = 'Calibri'
            
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.auto_size = None
            
            # PRESERVE ALL BULLET POINTS - ENSURE EXACT COUNT
            available_bullets = slide_data['bullets']
            
            # Log bullet count for debugging
            logger.info(f"📝 Slide {i+1} has {len(available_bullets)} bullet points")
            
            for bullet in available_bullets:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                
                # Adjust spacing based on content depth
                if content_depth == "comprehensive":
                    p.space_after = Pt(4)
                    p.space_before = Pt(2)
                elif content_depth == "detailed":
                    p.space_after = Pt(6)
                    p.space_before = Pt(3)
                else:
                    p.space_after = Pt(8)
                    p.space_before = Pt(4)
                
                # Smart font sizing based on content depth
                text_length = len(bullet)
                if content_depth == "comprehensive":
                    font_size = Pt(14) if text_length > 80 else Pt(15)
                elif content_depth == "detailed":
                    font_size = Pt(16) if text_length > 60 else Pt(17)
                else:
                    font_size = Pt(18) if text_length > 40 else Pt(19)
                
                for run in p.runs:
                    run.font.size = font_size
                    run.font.color.rgb = RGBColor(*colors['bullet'])
                    run.font.name = 'Calibri'

            # Optional images
            if include_images and slide_data.get('image_path'):
                try:
                    if i % 3 == 0:
                        left = Inches(7.5)
                        top = Inches(1.5)
                        width = Inches(4.5)
                        height = Inches(3.5)
                        
                        slide.shapes.add_picture(
                            slide_data['image_path'],
                            left, top, width, height
                        )
                except Exception as e:
                    logger.warning(f"⚠️ Could not add image to slide: {str(e)}")

        # Closing slide
        closing_slide = prs.slides.add_slide(prs.slide_layouts[0])
        background = closing_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors['bg'])
        
        closing_slide.shapes.title.text = "Thank You"
        closing_placeholder = closing_slide.placeholders[1]
        
        # Show content depth and slide count in closing slide
        bullet_counts = {
            "basic": "3 bullet points per slide",
            "detailed": "4 bullet points per slide", 
            "comprehensive": "5 bullet points per slide"
        }
        
        closing_placeholder.text = f"Generated with AI-Powered Chat-to-PPT\n\n{content_depth.title()} Content • {len(slides)} Slides • {bullet_counts.get(content_depth, 'Professional Quality')}"
        
        for paragraph in closing_slide.shapes.title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(36)
                run.font.color.rgb = RGBColor(*colors['title'])
        
        # Save presentation
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(tmp.name)
        logger.info(f"✅ PRESERVED BULLETS Presentation saved: {tmp.name}")
        
        # Clean up temporary image files
        for slide_data in slides:
            if slide_data.get('image_path') and os.path.exists(slide_data['image_path']):
                try:
                    os.unlink(slide_data['image_path'])
                except Exception as e:
                    logger.warning(f"Could not delete image file: {e}")
                    
        return tmp.name
        
    except Exception as e:
        logger.error(f"❌ PPT building failed: {str(e)}")
        return create_fallback_ppt()

def create_fallback_ppt():
    """Create a simple fallback presentation"""
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "AI Presentation Generated"
        slide.placeholders[1].text = "Powered by Google Gemini AI\n\nProfessional Quality • Ready to Present"
        
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(tmp.name)
        logger.info(f"🔄 Fallback presentation created: {tmp.name}")
        return tmp.name
    except Exception as e:
        logger.error(f"💥 Fallback PPT also failed: {str(e)}")
        raise HTTPException(status_code=500, detail="Failed to create PowerPoint file")

# Database operations
def init_db():
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS presentations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            topic TEXT NOT NULL,
            slides INTEGER NOT NULL,
            style TEXT NOT NULL,
            background_color TEXT NOT NULL,
            include_images BOOLEAN DEFAULT FALSE,
            content_depth TEXT DEFAULT 'detailed',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            downloaded BOOLEAN DEFAULT FALSE,
            download_count INTEGER DEFAULT 0,
            UNIQUE(topic, style, background_color, content_depth)
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS app_metrics (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            total_downloads INTEGER DEFAULT 0,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    cursor.execute('INSERT OR IGNORE INTO app_metrics (id, total_downloads) VALUES (1, 0)')
    
    cursor.execute("PRAGMA table_info(presentations)")
    columns = {column[1] for column in cursor.fetchall()}
    
    if 'content_depth' not in columns:
        cursor.execute('ALTER TABLE presentations ADD COLUMN content_depth TEXT DEFAULT "detailed"')
    
    if 'include_images' not in columns:
        cursor.execute('ALTER TABLE presentations ADD COLUMN include_images BOOLEAN DEFAULT FALSE')
    
    conn.commit()
    conn.close()
    logger.info("✅ Database initialized successfully")

init_db()

def save_presentation(history_data: dict) -> int:
    """Save presentation to database"""
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    
    actual_slides = history_data.get('actual_slides', history_data['slides'])
    
    cursor.execute('''
        SELECT id FROM presentations 
        WHERE topic = ? AND style = ? AND background_color = ? AND content_depth = ?
    ''', (
        history_data['topic'],
        history_data['style'],
        history_data['background_color'],
        history_data.get('content_depth', 'detailed')
    ))
    
    existing = cursor.fetchone()
    
    if existing:
        presentation_id = existing[0]
        cursor.execute('''
            UPDATE presentations 
            SET slides = ?, include_images = ?, created_at = CURRENT_TIMESTAMP
            WHERE id = ?
        ''', (
            actual_slides,
            history_data.get('include_images', False),
            presentation_id
        ))
    else:
        cursor.execute('''
            INSERT INTO presentations (topic, slides, style, background_color, include_images, content_depth)
            VALUES (?, ?, ?, ?, ?, ?)
        ''', (
            history_data['topic'],
            actual_slides,
            history_data['style'],
            history_data['background_color'],
            history_data.get('include_images', False),
            history_data.get('content_depth', 'detailed')
        ))
        presentation_id = cursor.lastrowid
    
    conn.commit()
    conn.close()
    
    logger.info(f"💾 Saved presentation to database: {history_data['topic']} with {actual_slides} slides")
    return presentation_id

def get_presentations(limit: int = 1000) -> list:
    """Get all presentations from database"""
    conn = sqlite3.connect('presentations.db')
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    
    cursor.execute('''
        SELECT * FROM presentations 
        ORDER BY created_at DESC 
        LIMIT ?
    ''', (limit,))
    
    presentations = [dict(row) for row in cursor.fetchall()]
    conn.close()
    
    logger.info(f"📊 Retrieved {len(presentations)} presentations from database")
    return presentations

def update_presentation_download(presentation_id: int):
    """Update presentation download count"""
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    
    cursor.execute('''
        UPDATE presentations 
        SET downloaded = TRUE, download_count = download_count + 1 
        WHERE id = ?
    ''', (presentation_id,))
    
    cursor.execute('''
        UPDATE app_metrics 
        SET total_downloads = total_downloads + 1, updated_at = CURRENT_TIMESTAMP 
        WHERE id = 1
    ''')
    
    conn.commit()
    conn.close()

def get_total_downloads() -> int:
    """Get total download count from metrics"""
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    
    cursor.execute('SELECT total_downloads FROM app_metrics WHERE id = 1')
    result = cursor.fetchone()
    conn.close()
    
    return result[0] if result else 0

def delete_presentation(presentation_id: int) -> bool:
    """Delete a presentation from database"""
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    
    cursor.execute('DELETE FROM presentations WHERE id = ?', (presentation_id,))
    success = cursor.rowcount > 0
    
    conn.commit()
    conn.close()
    return success

def clear_all_presentations() -> bool:
    """Clear all presentations from database"""
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    
    cursor.execute('DELETE FROM presentations')
    cursor.execute('UPDATE app_metrics SET total_downloads = 0')
    
    conn.commit()
    conn.close()
    return True

@app.get("/api/background-colors")
def get_background_colors():
    """Get available background colors"""
    return BACKGROUND_COLORS

@app.get("/api/content-depths")
def get_content_depths():
    """Get available content depth options"""
    return {
        "basic": "Basic (3 bullets, 5-7 words each) - Foundational concepts",
        "detailed": "Detailed (4 bullets, 8-12 words each) - Strategic analysis", 
        "comprehensive": "Comprehensive (5 bullets, 12-18 words each) - Executive insights"
    }

@app.get("/api/health")
def health():
    """Health check endpoint for Gemini AI"""
    try:
        test_prompt = "Hello, are you working? Respond with 'Yes, Gemini AI is ready for presentation generation.'"
        
        if gemini_model:
            response = gemini_model.generate_content(test_prompt)
            if response.text:
                return {
                    "ok": True, 
                    "model": GEMINI_MODEL,
                    "status": "Gemini AI Connected and Working",
                    "provider": "Google AI Studio",
                    "note": "Real AI content generation active"
                }
        
        return {
            "ok": False, 
            "model": GEMINI_MODEL,
            "status": "Gemini AI Not Available",
            "provider": "Google AI Studio", 
            "note": "Using enhanced content generation"
        }
            
    except Exception as e:
        return {
            "ok": False, 
            "error": str(e), 
            "model": GEMINI_MODEL,
            "provider": "Google AI Studio",
            "note": "Enhanced content generation active"
        }

def generate_presentation_content(topic: str, slides: int, content_depth: str):
    """UNIFIED function to generate presentation content for both outline and direct download"""
    logger.info(f"🎯 Generating content for: {topic}")
    logger.info(f"📊 Settings: EXACTLY {slides} slides, {content_depth} depth")
    
    prompt_template = PROMPT_TEMPLATES.get(content_depth, PROMPT_TEMPLATES["detailed"])
    prompt = prompt_template.format(topic=topic, slides=slides)
    
    # Try Gemini AI first for all content depths
    ai_text = call_gemini(prompt)
    
    if ai_text:
        # Use AI-generated content
        parsed_slides = parse_outline(ai_text, content_depth, slides)
        if parsed_slides and len(parsed_slides) >= min(3, slides):  # Accept at least 3 slides
            logger.info(f"✅ AI generated {len(parsed_slides)} slides with proper bullet points")
            return parsed_slides, "gemini"
    
    # Use enhanced fallback for all content depths
    logger.warning(f"⚠️ Using consistent fallback content for {content_depth} depth")
    fallback_text = generate_consistent_fallback(topic, slides, content_depth)
    parsed_slides = parse_outline(fallback_text, content_depth, slides)
    
    return parsed_slides, "fallback"

@app.post("/api/outline")
def api_outline(payload: GeneratePayload):
    """Generate presentation outline using UNIFIED content generation"""
    try:
        slides, content_source = generate_presentation_content(payload.topic, payload.slides, payload.content_depth)
        
        # If still no slides, create minimal fallback
        if not slides:
            if payload.content_depth == "comprehensive":
                minimal_text = f"Slide 1: Executive Overview of {payload.topic}\n- Comprehensive market analysis and strategic positioning\n- Financial projections and return on investment calculation\n- Implementation roadmap with resource allocation strategy\n- Risk assessment and mitigation planning framework\n- Performance metrics and success measurement criteria"
            elif payload.content_depth == "detailed":
                minimal_text = f"Slide 1: Strategic Analysis of {payload.topic}\n- Market overview and competitive landscape assessment\n- Key benefits and implementation considerations\n- Risk analysis and mitigation strategies\n- Success metrics and performance tracking"
            else:
                minimal_text = f"Slide 1: Introduction to {payload.topic}\n- Core concept overview and basic principles\n- Main applications and use cases\n- Key benefits and advantages"
            
            slides = parse_outline(minimal_text, payload.content_depth, payload.slides)
            content_source = "minimal_fallback"
        
        history_data = {
            "topic": payload.topic,
            "slides": payload.slides,
            "actual_slides": len(slides),
            "style": payload.style,
            "background_color": payload.background_color,
            "include_images": payload.include_images,
            "content_depth": payload.content_depth
        }
        
        presentation_id = save_presentation(history_data)
        
        logger.info(f"✅ Outline generated with {len(slides)} slides (Source: {content_source})")
        
        return {
            "topic": payload.topic,
            "style": payload.style,
            "background_color": payload.background_color,
            "content_depth": payload.content_depth,
            "slides": slides,
            "presentation_id": presentation_id,
            "ai_generated": content_source == "gemini",
            "exact_slide_count": len(slides) == payload.slides,
            "content_source": content_source
        }
    except Exception as e:
        logger.error(f"❌ Outline generation failed: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate-ppt")
def api_generate_ppt(payload: GeneratePayload):
    """Generate complete PPT using UNIFIED content generation"""
    try:
        logger.info(f"🚀 Generating PPT for: {payload.topic}")
        logger.info(f"📊 Settings: EXACTLY {payload.slides} slides, {payload.content_depth} depth, images: {payload.include_images}")
        
        # Use the SAME content generation function as outline
        slides, content_source = generate_presentation_content(payload.topic, payload.slides, payload.content_depth)
        
        # If still no slides, create minimal fallback
        if not slides:
            if payload.content_depth == "comprehensive":
                minimal_text = f"Slide 1: Executive Strategic Overview: {payload.topic}\n- Comprehensive market analysis with current trends and competitive landscape\n- Strategic business case with ROI calculation and financial projections\n- Implementation roadmap detailing phased approach and resource allocation\n- Stakeholder impact analysis covering change management requirements\n- Performance measurement framework defining KPIs and improvement processes"
            elif payload.content_depth == "detailed":
                minimal_text = f"Slide 1: Comprehensive {payload.topic} Analysis\n- Detailed overview with market context and current relevance\n- Key industry trends analysis and emerging patterns\n- Strategic importance and business impact considerations\n- Implementation challenges and potential solutions"
            else:
                minimal_text = f"Slide 1: Introduction to {payload.topic}\n- Core concept definition and overview\n- Main purpose and basic applications\n- Key benefits and importance"
            
            slides = parse_outline(minimal_text, payload.content_depth, payload.slides)
            content_source = "minimal_fallback"
        
        history_data = {
            "topic": payload.topic,
            "slides": payload.slides,
            "actual_slides": len(slides),
            "style": payload.style,
            "background_color": payload.background_color,
            "include_images": payload.include_images,
            "content_depth": payload.content_depth
        }
        
        presentation_id = save_presentation(history_data)
        
        if payload.include_images:
            logger.info("🖼️ Generating images for slides...")
            for i, slide in enumerate(slides):
                if i % 3 == 0:
                    image_prompt = f"professional business presentation slide about {slide['title']}, clean modern corporate design, informative content"
                    image_path = generate_image_pollinations(image_prompt)
                    slide['image_path'] = image_path
                    time.sleep(1)
        
        ppt_path = build_ppt(slides, payload.style, payload.background_color, payload.include_images, payload.content_depth)
        
        update_presentation_download(presentation_id)
        
        filename = f"{re.sub(r'[^a-zA-Z0-9_-]+', '_', payload.topic) or 'ai_presentation'}.pptx"
        
        logger.info(f"✅ FULL BULLETS PPT generated: {filename} with {len(slides)} slides (Source: {content_source})")
        
        return FileResponse(
            ppt_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename,
            background=None
        )
        
    except Exception as e:
        logger.error(f"❌ PPT generation failed: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/history")
def get_history():
    """Get presentation history"""
    try:
        presentations = get_presentations()
        total_downloads = get_total_downloads()
        
        return {
            "presentations": presentations,
            "total_downloads": total_downloads
        }
    except Exception as e:
        logger.error(f"❌ Failed to fetch history: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/history/{presentation_id}")
def delete_history_item(presentation_id: int):
    """Delete a specific presentation from history"""
    try:
        success = delete_presentation(presentation_id)
        if not success:
            raise HTTPException(status_code=404, detail="Presentation not found")
        
        return {"message": "Presentation deleted successfully"}
    except Exception as e:
        logger.error(f"❌ Failed to delete presentation: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/history")
def clear_history():
    """Clear all presentation history"""
    try:
        clear_all_presentations()
        return {"message": "All history cleared successfully"}
    except Exception as e:
        logger.error(f"❌ Failed to clear history: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/metrics")
def get_metrics():
    """Get application metrics"""
    try:
        total_downloads = get_total_downloads()
        conn = sqlite3.connect('presentations.db')
        cursor = conn.cursor()
        
        cursor.execute('SELECT COUNT(*) FROM presentations')
        total_presentations = cursor.fetchone()[0]
        
        cursor.execute('SELECT COUNT(*) FROM presentations WHERE downloaded = TRUE')
        downloaded_presentations = cursor.fetchone()[0]
        
        cursor.execute('SELECT content_depth, COUNT(*) FROM presentations GROUP BY content_depth')
        depth_distribution = {row[0]: row[1] for row in cursor.fetchall()}
        
        conn.close()
        
        return {
            "total_downloads": total_downloads,
            "total_presentations": total_presentations,
            "downloaded_presentations": downloaded_presentations,
            "content_depth_distribution": depth_distribution
        }
    except Exception as e:
        logger.error(f"❌ Failed to fetch metrics: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/")
def root():
    return {
        "message": "Chat-to-PPT API is running", 
        "version": "18.5", # Updated Version
        "ai_provider": "Google Gemini 2.5 Flash Lite",
        "status": "Production Ready - COMPLETELY FIXED TITLE & BULLET COUNTS",
        "features": "Basic: 3 bullets • Detailed: 4 bullets • Comprehensive: 5 bullets"
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
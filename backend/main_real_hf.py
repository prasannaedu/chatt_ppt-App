from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
import os, time, requests, re, tempfile, sqlite3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Hugging Face Configuration
HUGGINGFACE_API_KEY = os.getenv("HUGGINGFACE_API_KEY")
HUGGINGFACE_MODEL = os.getenv("HUGGINGFACE_MODEL", "google/flan-t5-base")
HF_API_URL = "https://router.huggingface.co/hf-inference"

app = FastAPI(title="Chat-to-PPT API - Real AI", version="8.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

class GeneratePayload(BaseModel):
    topic: str
    slides: int = 5
    style: str = "Blue-Professional"
    background_color: str = "#FFFFFF"
    include_images: bool = False
    content_depth: str = "detailed"

# Enhanced prompts for better AI responses
PROMPT_TEMPLATES = {
    "basic": "Create a presentation outline about '{topic}' with exactly {slides} slides. Each slide should have a title and 3 concise bullet points. Format: Slide 1: [Title] - [Bullet 1] - [Bullet 2] - [Bullet 3]",
    
    "detailed": "Generate a detailed professional presentation outline about '{topic}' with exactly {slides} slides. Each slide should have an engaging title and 4 informative bullet points with explanations. Format: Slide 1: [Title] - [Detailed bullet 1] - [Detailed bullet 2] - [Detailed bullet 3] - [Detailed bullet 4]",
    
    "comprehensive": "Develop a comprehensive executive presentation outline about '{topic}' with exactly {slides} slides. Each slide should have a strategic title and 5 in-depth bullet points with data, insights, and recommendations. Format: Slide 1: [Title] - [Strategic insight 1] - [Data point 2] - [Analysis 3] - [Recommendation 4] - [Action item 5]"
}

def call_real_huggingface(prompt: str, max_retries: int = 3) -> str:
    """Call REAL Hugging Face API with new router endpoint"""
    if not HUGGINGFACE_API_KEY:
        raise Exception("Hugging Face API key not configured")
    
    headers = {
        "Authorization": f"Bearer {HUGGINGFACE_API_KEY}",
        "Content-Type": "application/json"
    }
    
    payload = {
        "inputs": prompt,
        "model": HUGGINGFACE_MODEL,
        "parameters": {
            "max_new_tokens": 800,
            "temperature": 0.7,
            "do_sample": True,
            "top_p": 0.9
        }
    }
    
    for attempt in range(max_retries):
        try:
            logger.info(f"🤖 Calling REAL Hugging Face API (attempt {attempt + 1})")
            logger.info(f"📝 Prompt: {prompt[:100]}...")
            
            response = requests.post(
                HF_API_URL,
                headers=headers,
                json=payload,
                timeout=60
            )
            
            logger.info(f"🔧 API Response Status: {response.status_code}")
            
            if response.status_code == 200:
                result = response.json()
                logger.info("✅ REAL AI Response received!")
                
                # Extract generated text
                if isinstance(result, list) and len(result) > 0:
                    generated_text = result[0].get('generated_text', '')
                elif isinstance(result, dict):
                    generated_text = result.get('generated_text', str(result))
                else:
                    generated_text = str(result)
                
                if generated_text and len(generated_text) > 20:
                    logger.info(f"📄 Real AI Content Length: {len(generated_text)}")
                    logger.info(f"📄 Content Preview: {generated_text[:200]}...")
                    return generated_text
                else:
                    logger.warning("⚠️ Empty response from real AI")
                    continue
                    
            elif response.status_code == 503:
                wait_time = 30
                logger.info(f"⏳ Real AI model loading, waiting {wait_time} seconds...")
                time.sleep(wait_time)
                continue
                
            else:
                logger.error(f"❌ Real AI API error {response.status_code}: {response.text}")
                if attempt < max_retries - 1:
                    time.sleep(15)
                    continue
                else:
                    raise Exception(f"Real AI API failed: {response.status_code}")
                    
        except requests.exceptions.Timeout:
            logger.error("⏰ Real AI API timeout")
            if attempt < max_retries - 1:
                time.sleep(15)
                continue
            else:
                raise Exception("Real AI API timeout after retries")
                
        except Exception as e:
            logger.error(f"💥 Real AI API exception: {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(10)
                continue
            else:
                raise
    
    raise Exception("All real AI API attempts failed")

def parse_outline(text: str, content_depth: str = "detailed") -> List[dict]:
    """Parse AI-generated outline into structured slides"""
    slides = []
    blocks = re.split(r"\n\s*\n", text.strip())
    
    max_bullets = {"basic": 3, "detailed": 4, "comprehensive": 5}.get(content_depth, 4)
    
    for block in blocks:
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        if not lines:
            continue
            
        # Extract title (handle different formats)
        title_line = lines[0]
        title = re.sub(r"^Slide\s*\d+:\s*", "", title_line, flags=re.I).strip()
        if not title or len(title) < 2:
            title = "AI Generated Slide"
            
        # Extract bullets
        bullets = []
        for ln in lines[1:]:
            if ln.startswith('-'):
                clean_line = ln[1:].strip()
                if clean_line and len(clean_line) > 3:
                    bullets.append(clean_line)
        
        bullets = bullets[:max_bullets]
        slides.append({"title": title, "bullets": bullets})
    
    # Fallback if parsing fails
    if not slides:
        slides = [{"title": "AI Presentation", "bullets": ["Generated with Real AI", "Professional Content"]}]
    
    return slides

def build_ppt(slides, style: str, background_color: str = "#FFFFFF") -> str:
    """Build PowerPoint from AI-generated content"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        if slides and len(slides) > 0:
            title_shape.text = slides[0]['title']
            if slides[0].get('bullets'):
                subtitle_text = "\n".join([f"• {bullet}" for bullet in slides[0]['bullets'][:3]])
                subtitle_shape.text = subtitle_text
        else:
            title_shape.text = "AI-Generated Presentation"
            subtitle_shape.text = "Powered by Real Hugging Face AI"
        
        # Format title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(44)
                run.font.bold = True
        
        # Content slides
        for slide_data in slides[1:] if len(slides) > 1 else slides:
            content_slide = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide)
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True
            
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for bullet in slide_data['bullets']:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = Pt(12)
                
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.name = 'Calibri'
        
        # Save presentation
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(tmp.name)
        logger.info(f"✅ Real AI Presentation saved: {tmp.name}")
        return tmp.name
        
    except Exception as e:
        logger.error(f"❌ PPT building failed: {str(e)}")
        raise

# Database setup
def init_db():
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS presentations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            topic TEXT NOT NULL,
            slides INTEGER NOT NULL,
            style TEXT NOT NULL,
            background_color TEXT NOT NULL,
            content_depth TEXT DEFAULT 'detailed',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.commit()
    conn.close()
    logger.info("✅ Database initialized")

init_db()

def save_presentation(history_data: dict) -> int:
    conn = sqlite3.connect('presentations.db')
    cursor = conn.cursor()
    cursor.execute('''
        INSERT INTO presentations (topic, slides, style, background_color, content_depth)
        VALUES (?, ?, ?, ?, ?)
    ''', (
        history_data['topic'], history_data['slides'], history_data['style'],
        history_data['background_color'], history_data.get('content_depth', 'detailed')
    ))
    presentation_id = cursor.lastrowid
    conn.commit()
    conn.close()
    return presentation_id

# API Endpoints
@app.get("/api/health")
def health():
    """Health check with real AI test"""
    try:
        # Test with real AI
        test_response = call_real_huggingface("Say 'AI is working' in one sentence.")
        return {
            "ok": True,
            "model": HUGGINGFACE_MODEL,
            "status": "Real AI Connected",
            "ai_response": test_response[:100],
            "provider": "Hugging Face Real AI"
        }
    except Exception as e:
        return {
            "ok": False,
            "error": str(e),
            "model": HUGGINGFACE_MODEL,
            "status": "AI Connection Failed"
        }

@app.post("/api/generate-ppt")
def api_generate_ppt(payload: GeneratePayload):
    """Generate PPT with REAL AI content"""
    try:
        logger.info(f"🚀 Generating PPT with REAL AI for: {payload.topic}")
        logger.info(f"📊 Settings: {payload.slides} slides, {payload.content_depth} depth")
        
        # Get prompt template
        prompt_template = PROMPT_TEMPLATES.get(payload.content_depth, PROMPT_TEMPLATES["detailed"])
        prompt = prompt_template.format(topic=payload.topic, slides=payload.slides)
        
        # Call REAL Hugging Face AI
        logger.info("🤖 Calling REAL Hugging Face AI...")
        ai_content = call_real_huggingface(prompt)
        
        # Parse AI-generated content
        slides = parse_outline(ai_content, payload.content_depth)
        logger.info(f"📄 Real AI generated {len(slides)} slides")
        
        # Save to database
        history_data = {
            "topic": payload.topic,
            "slides": payload.slides,
            "style": payload.style,
            "background_color": payload.background_color,
            "content_depth": payload.content_depth
        }
        
        presentation_id = save_presentation(history_data)
        
        # Build PowerPoint
        ppt_path = build_ppt(slides, payload.style, payload.background_color)
        
        filename = f"{payload.topic.replace(' ', '_')}_real_ai.pptx"
        
        logger.info(f"✅ REAL AI PPT generated successfully: {filename}")
        
        return FileResponse(
            ppt_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename
        )
        
    except Exception as e:
        logger.error(f"❌ REAL AI PPT generation failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Real AI Error: {str(e)}")

@app.get("/")
def root():
    return {
        "message": "Chat-to-PPT API with REAL AI is running!",
        "version": "8.0",
        "ai_provider": "Hugging Face Real AI",
        "model": HUGGINGFACE_MODEL,
        "status": "Ready for Real AI Presentations"
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
